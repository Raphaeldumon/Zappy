#!/usr/bin/env python3
"""Assemble Keynote.pptx : fonds PNG plein écran + zones de texte natives.

Source de vérité des zones : le manifeste SLIDES ci-dessous. Les rectangles
"px" doivent rester identiques aux divs .native des slideNN.html.
Conversions : 1 px design = 6350 EMU ; pt natif = px / 2.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

HERE = Path(__file__).parent
OUT = HERE / "out"
TARGET = HERE.parent / "Keynote.pptx"

EMU_PER_PX = 6350          # 12192000 EMU / 1920 px
SLIDE_W = 12192000         # 16:9
SLIDE_H = 6858000

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

SLIDES = [
    {
        "n": "01",
        "zones": [
            {"px": (310, 640, 1300, 60), "pt": 20, "align": "center",
             "text": "Une course d’intelligences artificielles sur un monde-anneau"},
            {"px": (210, 950, 1500, 44), "pt": 14, "align": "center", "color": "8E9BC0",
             "text": "Raphael Dumon · Clément Fabre · Mathéo Emma · Léo Cazabat · Joachim Ismael"},
        ],
    },
    {
        "n": "02",
        "zones": [
            {"px": (110, 400, 780, 300), "pt": 20,
             "text": "Un serveur arbitre la partie et détient la seule vérité du monde.\n"
                     "Des drones autonomes, chacun piloté par une IA.\n"
                     "Une interface 3D qui montre tout, en direct."},
            {"px": (110, 760, 780, 90), "pt": 16, "color": "8E9BC0",
             "text": "Trois programmes séparés, écrits par nous, qui ne se parlent "
                     "qu’en s’envoyant du texte sur le réseau."},
        ],
    },
    {
        "n": "03",
        "zones": [
            {"px": (110, 400, 800, 340), "pt": 20,
             "text": "Manger pour survivre — la nourriture s’épuise en continu.\n"
                     "Collecter six types de pierres précieuses.\n"
                     "S’élever par des incantations — parfois à plusieurs."},
            {"px": (110, 790, 800, 120), "pt": 22, "bold": True, "color": "F5C866",
             "text": "Victoire : la première équipe qui amène 6 joueurs au niveau 8."},
        ],
    },
    {
        "n": "04",
        "zones": [
            {"px": (110, 400, 800, 300), "pt": 20,
             "text": "Sortir à droite, c’est revenir par la gauche : la carte est un tore.\n"
                     "Le temps s’écoule en ticks : chaque action a un prix.\n"
                     "La fréquence f accélère le monde entier — de la balade au sprint."},
            {"px": (110, 760, 800, 90), "pt": 16, "color": "8E9BC0",
             "text": "f = 1 → une unité de temps par seconde. f = 1000 → mille fois plus vite."},
        ],
    },
    {
        "n": "05",
        "zones": [
            {"px": (110, 400, 800, 300), "pt": 20,
             "text": "Les IA et la GUI sont les clients : ils passent commande.\n"
                     "Le serveur vérifie que c’est au menu, répond ok ou ko…\n"
                     "…puis sert chaque plat au bon moment — pas avant."},
            {"px": (110, 760, 800, 90), "pt": 16, "color": "8E9BC0",
             "text": "Hors menu ? « ko », et on passe à la commande suivante."},
        ],
    },
    {
        "n": "06",
        "zones": [
            {"px": (110, 400, 800, 320), "pt": 20,
             "text": "Un seul programme, un seul fil d’exécution — et aucune attente active.\n"
                     "poll() surveille toutes les tables à la fois.\n"
                     "Un agenda d’événements planifie chaque échéance : le serveur dort, "
                     "puis se réveille exactement quand il faut."},
            {"px": (110, 790, 800, 90), "pt": 16, "color": "8E9BC0",
             "text": "Réveillé par un paquet réseau ou par l’échéance suivante — jamais pour rien."},
        ],
    },
    {
        "n": "07",
        "zones": [
            {"px": (110, 400, 820, 320), "pt": 20,
             "text": "Il ne sait pas où il est — aucune commande ne donne sa position.\n"
                     "Il ne sait pas combien ils sont — pas d’annuaire d’équipe.\n"
                     "Il ne sait pas qui parle — un message entendu est anonyme."},
            {"px": (110, 790, 820, 90), "pt": 18, "bold": True,
             "text": "Tout ce qui suit existe pour contourner ces trois handicaps."},
        ],
    },
    {
        "n": "08",
        "zones": [
            {"px": (110, 400, 820, 320), "pt": 20,
             "text": "Chaque son entendu arrive avec une direction — huit possibles.\n"
                     "HELLO : chacun se présente en boucle, l’équipe se compte toute seule.\n"
                     "Le plus petit identifiant devient le chef — élection sans dispute."},
            {"px": (110, 790, 820, 90), "pt": 16, "color": "8E9BC0",
             "text": "Pour rejoindre quelqu’un : marcher vers son son, encore et encore."},
        ],
    },
    {
        "n": "09",
        "zones": [
            {"px": (110, 400, 820, 320), "pt": 17,
             "text": "Vu d’un chef, un « joueur » sur sa case n’a pas de niveau : "
                     "impossible de savoir si c’est le bon coéquipier.\n"
                     "Mais un son reçu « sur ma case » (direction 0), lui, ne ment pas.\n"
                     "Chacun annonce ARRIVED en entendant le chef à direction 0 — "
                     "le rituel ne part que sur les présences confirmées."},
            {"px": (110, 820, 820, 80), "pt": 16, "color": "8E9BC0",
             "text": "Une idée simple, trouvée après beaucoup de rituels ratés."},
        ],
    },
    {
        "n": "10",
        "zones": [
            {"px": (110, 400, 1100, 220), "pt": 20,
             "text": "Une seule stratégie, réglée comme une horlogerie : chacun amasse "
                     "tout le nécessaire, l’équipe fait le plein de nourriture, six "
                     "convergent vers le chef…\n"
                     "…qui dépose les 38 pierres et enchaîne les six rituels d’un coup."},
            {"px": (110, 950, 900, 70), "pt": 18, "bold": True, "color": "5CE8B5",
             "text": "Du niveau 2 au niveau 8 : environ 3 secondes. Partie gagnée en ~30."},
        ],
    },
    {
        "n": "11",
        "zones": [
            {"px": (110, 400, 800, 300), "pt": 20,
             "text": "Au branchement : le serveur envoie un instantané complet du monde.\n"
                     "Ensuite : chaque changement arrive en temps réel, message par message…\n"
                     "…et la 3D le traduit aussitôt — déplacements, rituels, naissances, morts."},
            {"px": (110, 760, 800, 90), "pt": 16, "color": "8E9BC0",
             "text": "La GUI observe tout mais ne joue jamais : spectatrice, pas arbitre."},
        ],
    },
    {
        "n": "13",
        "zones": [
            {"px": (110, 400, 900, 300), "pt": 20,
             "text": "Chaque message reçu est daté et archivé.\n"
                     "Tirer la timeline en arrière reconstruit le monde à cet instant "
                     "précis — en rejouant son histoire.\n"
                     "Une partie entière s’enregistre dans un fichier .zrec, et se revoit plus tard."},
            {"px": (110, 760, 900, 80), "pt": 16, "color": "8E9BC0",
             "text": "Idéal pour comprendre une défaite… ou savourer une victoire image par image."},
        ],
    },
]


def add_zone(slide, zone):
    x, y, w, h = (Emu(v * EMU_PER_PX) for v in zone["px"])
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = zone["text"].split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ALIGN[zone.get("align", "left")]
        if i > 0:
            p.space_before = Pt(zone.get("gap_pt", 8))
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = "Arial"
        f.size = Pt(zone["pt"])
        f.bold = zone.get("bold", False)
        f.color.rgb = RGBColor.from_string(zone.get("color", "EAF0FF"))


def main():
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    blank = prs.slide_layouts[6]
    for spec in SLIDES:
        png = OUT / f"slide{spec['n']}.png"
        if not png.exists():
            raise SystemExit(f"fond manquant : {png} (lancer render.sh)")
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=Emu(SLIDE_W), height=Emu(SLIDE_H))
        for zone in spec["zones"]:
            add_zone(slide, zone)
    prs.save(TARGET)
    print(f"écrit : {TARGET} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
